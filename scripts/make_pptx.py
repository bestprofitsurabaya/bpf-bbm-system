#!/usr/bin/env python3
"""Buat presentasi PPTX (16:9) dari materi PRESENTASI.md — bisa diedit di PowerPoint.

Cara pakai:
    python3 scripts/make_pptx.py
Hasil: presentasi/BPF_WorkHub_Presentasi.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
SOFT = RGBColor(0xEF, 0xF6, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)
BADGE = "presentasi/bpf-badge.png"  # badge logo bulat BPF (dibuat oleh make_videos.sh / generate badge)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

import os


def logo(slide, size=Inches(0.5), x=Inches(0.3), y=Inches(6.95)):
    """Tempel badge logo BPF (bulat, cincin biru) — aman bila file tidak ada."""
    if os.path.isfile(BADGE):
        try:
            slide.shapes.add_picture(BADGE, x, y, width=size, height=size)
        except Exception as e:
            print(f"⚠ logo skip: {e}")


def new_slide():
    return prs.slides.add_slide(BLANK)


def box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill if fill is not None else WHITE
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def text(slide, x, y, w, h, runs, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, t in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb


def kicker(slide, label):
    text(slide, Inches(0.9), Inches(0.45), Inches(11), Inches(0.4), label.upper(),
         size=13, color=BLUE, bold=True)


def title(slide, label):
    text(slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.9), label,
         size=34, color=DARK, bold=True)


def keybox(slide, label, y=Inches(5.9)):
    box(slide, Inches(0.9), y, Inches(11.5), Inches(0.95), fill=SOFT)
    text(slide, Inches(1.15), y + Inches(0.18), Inches(11), Inches(0.7),
         label, size=16, color=DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, items, y=Inches(1.9), size=16, color=GRAY, spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.2), Inches(3.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = spacing
        r = p.add_run()
        r.text = "→  " + t
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def cards(slide, items, y=Inches(2.0)):
    n = len(items)
    gap = Inches(0.3)
    w = Inches((11.5 - gap * (n - 1)) / n)
    h = Inches(3.2)
    for i, (ico, head, body) in enumerate(items):
        x = Inches(0.9) + Emu(int((w + gap) * i))
        box(slide, x, y, w, h, fill=WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
        text(slide, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.6),
             ico, size=28)
        text(slide, x + Inches(0.25), y + Inches(1.0), w - Inches(0.5), Inches(0.6),
             head, size=17, color=DARK, bold=True)
        text(slide, x + Inches(0.25), y + Inches(1.55), w - Inches(0.5), Inches(1.5),
             body, size=13, color=GRAY)


def table(slide, rows, y=Inches(1.95), h=Inches(3.4), col_w=None, header=True):
    from pptx.util import Inches as In
    nrows, ncols = len(rows), len(rows[0])
    gt = slide.shapes.add_table(nrows, ncols, Inches(0.9), y, Inches(11.5), h).table
    if col_w:
        for c, w in enumerate(col_w):
            gt.columns[c].width = In(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = In(0.08)
            cell.margin_right = In(0.08)
            cell.margin_top = In(0.03)
            cell.margin_bottom = In(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.0
            r = p.add_run()
            r.text = val
            if header and ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SOFT
                r.font.size = Pt(13)
                r.font.bold = True
                r.font.color.rgb = BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
                r.font.size = Pt(14)
                r.font.bold = (ci == 0)
                r.font.color.rgb = DARK if ci == 0 else GRAY


def footer(slide, n, total):
    logo(slide)  # badge kecil kiri bawah
    text(slide, Inches(11.9), Inches(7.05), Inches(1.2), Inches(0.35),
         f"{n} / {total}", size=11, color=GRAY, align=PP_ALIGN.RIGHT)


# ================================================================
# SLIDE 1 — Judul
# ================================================================
s = new_slide()
box(s, 0, 0, SW, SH, fill=LIGHT)
box(s, 0, 0, SW, Inches(0.12), fill=BLUE)
kicker(s, "Materi Presentasi · Demo Aplikasi")
logo(s, size=Inches(1.7), x=Inches(10.9), y=Inches(0.35))  # badge besar kanan atas
text(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.9),
     ["BPF WorkHub"], size=48, color=DARK, bold=True)
text(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(0.6),
     "Satu aplikasi untuk semua urusan armada, BBM, kasbon, log perjalanan, appointment, dan air minum.",
     size=20, color=GRAY)
cards(s, [
    ("🚚", "Driver", "Klaim BBM · Kasbon · Trip · Offline"),
    ("💧", "OB", "Pengajuan air minum + foto bukti"),
    ("🧾", "GA & Finance", "Verifikasi, persetujuan, rekap"),
    ("📣", "Marketing & Chief Driver", "Appointment & pembagian tugas"),
], y=Inches(3.8))
footer(s, 1, 12)

# ================================================================
# SLIDE 2 — Cerita Masalah
# ================================================================
s = new_slide()
kicker(s, "Pembuka")
title(s, "Cerita dari Lapangan")
bullets(s, [
    "Sopir bayar BBM pakai uang sendiri — lalu menunggu berhari-hari uangnya diganti.",
    "Finance merekap struk manual di Excel — sering bingung \u201cuang ini untuk yang mana?\u201d.",
    "OB beli galon air minum — tidak ada bukti resmi kapan, berapa, dari siapa.",
    "Laporan perjalanan hilang atau telat, tidak ada jejak siapa yang mengerjakan apa.",
], y=Inches(1.95), size=19, color=GRAY, spacing=1.6)
keybox(s, "💡  Kunci: dari \u201ccatat di kertas, rekap di Excel\u201d menjadi satu aplikasi, semua tercatat otomatis, semua bisa diaudit.")
footer(s, 2, 12)

# ================================================================
# SLIDE 3 — Satu Aplikasi Semua Peran
# ================================================================
s = new_slide()
kicker(s, "Gambaran Besar")
title(s, "Satu Aplikasi, Semua Peran")
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
     "Setiap orang login dengan Username + PIN 6 digit dan langsung masuk ke halamannya masing-masing.",
     size=18, color=GRAY)
table(s, [
    ["Peran", "Langsung dibawa ke"],
    ["Admin", "Dashboard Admin — seluruh operasi + manajemen user"],
    ["GA", "Dashboard GA — klaim, kasbon, trip, verifikasi anomali"],
    ["Finance", "Dashboard Finance — rekap air minum & kasbon"],
    ["Driver", "Aplikasi Driver (PWA) — BBM, kasbon, trip, rapor"],
    ["OB", "Halaman Air Minum"],
    ["Marketing / Chief Driver", "Marketing Hub / Board Chief Driver"],
], y=Inches(2.35), h=Inches(3.2), col_w=[4.2, 7.3])
keybox(s, "🔒  Kunci: satu aplikasi untuk semua — tapi setiap orang hanya melihat dunia yang menjadi tugasnya (least privilege).")
footer(s, 3, 12)

# ================================================================
# SLIDE 4 — Air Minum
# ================================================================
s = new_slide()
kicker(s, "Demo Utama 1")
title(s, "Air Minum: OB → Finance → PDF")
cards(s, [
    ("📝", "1. OB mengisi", "Tanggal, jumlah, jenis (Gelas/Botol/Galon), merk + foto sebelum & sesudah"),
    ("✅", "2. Finance verifikasi", "Periksa bukti, isi remark & catatan — semua terekam"),
    ("📄", "3. PDF tanda terima", "Ditandatangani Finance (penyerah) & GA (penerima)"),
], y=Inches(2.0))
bullets(s, [
    "Nama OB asli: Faisol, Febri, Edwin — masing-masing hanya melihat pengajuannya sendiri.",
    "Tanpa foto bukti, pengajuan tidak bisa diproses.",
], y=Inches(4.55), size=16)
keybox(s, "🏁  Kunci: dari pembelian galon sampai tanda terima resmi — tercatat, terbukti, bisa diaudit, tanpa kertas.")
footer(s, 4, 12)

# ================================================================
# SLIDE 5 — Dashboard per Peran
# ================================================================
s = new_slide()
kicker(s, "Demo Utama 2")
title(s, "Dashboard per Peran")
table(s, [
    ["Halaman", "Yang bisa dilakukan langsung"],
    ["Dashboard GA", "✅ Approve · ✕ Tolak klaim · 🛡 Verifikasi anomali · kasbon DRAFT · trip pending"],
    ["Dashboard Finance", "Rekap air minum per OB / jenis / merk · antrean verifikasi · kasbon · ⬇ Export CSV"],
    ["Dashboard Admin", "Ringkasan seluruh operasi + Users, Settings, Audit Log"],
], y=Inches(1.95), h=Inches(2.9), col_w=[3.4, 8.1])
bullets(s, [
    "Alasan penolakan wajib diisi — untuk jejak audit.",
    "Klaim bertanda ⚠️ melewati verifikasi anomali khusus.",
], y=Inches(4.7), size=16)
keybox(s, "🎯  Kunci: pekerjaan selesai dari satu halaman — tidak berpindah-pindah menu.")
footer(s, 5, 12)

# ================================================================
# SLIDE 6 — Kasbon & Kode Unik
# ================================================================
s = new_slide()
kicker(s, "Demo Utama 3")
title(s, "Klaim BBM & Kasbon")
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5),
     "Alur relay: Driver → GA → Finance → LPJ → Selesai 🎉", size=18, color=GRAY)
table(s, [
    ["Tahap", "Siapa", "Status"],
    ["Ajukan klaim + foto struk", "Driver", "Draft / Pending"],
    ["Setujui & serahkan dana", "GA", "Disetujui GA"],
    ["Cairkan dana", "Finance", "Diserahkan"],
    ["Isi LPJ (pertanggungjawaban)", "Driver", "LPJ Diajukan"],
    ["Verifikasi LPJ", "GA", "Selesai 🎉"],
], y=Inches(2.35), h=Inches(3.1), col_w=[5.6, 2.9, 3.0])
keybox(s, "🔢  Kode unik: kasbon Rp 100.000 otomatis menjadi Rp 100.023 — nominal persis inilah bukti \u201cuang ini untuk kasbon siapa\u201d.")
footer(s, 6, 12)

# ================================================================
# SLIDE 7 — Realtime
# ================================================================
s = new_slide()
kicker(s, "Demo Utama 4")
title(s, "Realtime & Notifikasi")
cards(s, [
    ("🔔", "Toast & Lonceng", "Pengajuan baru muncul seketika — ikon per jenis (klaim, trip, air minum, appointment)"),
    ("🔄", "Auto-refresh", "Antrean GA/Finance ter-refresh sendiri tanpa muat ulang"),
    ("🟢", "Indikator hidup", "\u201c⚡ Realtime\u201d di pojok — jika putus tampil merah"),
], y=Inches(2.1))
keybox(s, "📡  Kunci: aplikasi ini \u201chidup\u201d — begitu ada yang mengajukan, yang berwenang langsung tahu, bukan besok pagi.")
footer(s, 7, 12)

# ================================================================
# SLIDE 8 — Driver Offline
# ================================================================
s = new_slide()
kicker(s, "Demo Utama 5")
title(s, "Driver: Offline di Jalan")
bullets(s, [
    "Login PIN — identitas menempel otomatis di semua laporan (tidak bisa mengaku-ngaku).",
    "4 tab: ⛽ BBM · 💰 Kasbon · 🗺️ Trip · 📊 Rapor.",
    "Offline-first: data tersimpan di HP, terkirim otomatis saat sinyal kembali + tombol 🔄 Sinkron.",
    "Rapor performa: masukkan nopol → HEMAT / CUKUP / BOROS.",
    "Bisa dipasang di layar utama HP seperti aplikasi biasa (PWA).",
], y=Inches(2.0), size=19, spacing=1.7)
keybox(s, "🌐  Kunci: sopir di jalan tetap produktif — sinyal hilang bukan alasan data hilang.")
footer(s, 8, 12)

# ================================================================
# SLIDE 9 — Marketing & Chief Driver
# ================================================================
s = new_slide()
kicker(s, "Demo Pendukung")
title(s, "Marketing & Chief Driver")
cards(s, [
    ("📅", "Marketing Hub", "Input appointment, pantau status, lihat hasil kunjungan — data penjualan tersaji otomatis"),
    ("🚛", "Chief Driver", "Papan pembagian tugas: kendaraan belum ditugaskan, beban tiap driver, unduh rekap harian Excel"),
], y=Inches(2.2))
keybox(s, "🗂️  Kunci: bukan hanya BBM — penjadwalan kunjungan & pembagian sopir juga satu pintu.")
footer(s, 9, 12)

# ================================================================
# SLIDE 10 — Keamanan
# ================================================================
s = new_slide()
kicker(s, "Kepercayaan")
title(s, "Keamanan & Tata Kelola")
table(s, [
    ["Aspek", "Jawaban singkat"],
    ["Login PIN per orang", "Tidak ada akun bersama — semua aksi tercatat siapa"],
    ["Role-based access", "GA tak bisa buka rekap finance; OB tak bisa lihat klaim"],
    ["Audit Log", "Semua aksi penting tercatat & bisa difilter"],
    ["Anti peretasan", "Percobaan login dibatasi · token CSRF · unggahan diperiksa"],
    ["Kerahasiaan", "Sesi aman (HTTPS) · cookie ketat · foto bukti terlindungi"],
], y=Inches(1.95), h=Inches(3.4), col_w=[4.0, 7.5])
keybox(s, "🛡️  Kunci: keamanan berlapis — dari PIN pengguna sampai jejak audit setiap transaksi.")
footer(s, 10, 12)

# ================================================================
# SLIDE 11 — Kualitas
# ================================================================
s = new_slide()
kicker(s, "Kredibilitas")
title(s, "Kualitas & Kepatuhan")
cards(s, [
    ("✅", "179 uji otomatis", "97 backend + 82 frontend — setiap perubahan diuji"),
    ("🌐", "E2E di produksi", "Alur kritis diverifikasi di lingkungan nyata"),
    ("📜", "Standar ISO", "27001 (keamanan) · 9241-11 (UX) · 9001 (mutu)"),
    ("♿", "Aksesibilitas", "Fokus keyboard, kontras ≥ 4,5:1, mode kontras tinggi & gelap"),
], y=Inches(2.1))
keybox(s, "🏆  Kunci: bukan prototipe — sistem yang diuji, diamankan, dan siap dipakai harian.")
footer(s, 11, 12)

# ================================================================
# SLIDE 12 — Penutup
# ================================================================
s = new_slide()
box(s, 0, 0, SW, SH, fill=LIGHT)
box(s, 0, 0, SW, Inches(0.12), fill=BLUE)
kicker(s, "Penutup")
title(s, "Rangkuman & Langkah Berikutnya")
text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.4),
     ["1  ·  0  ·  100%"], size=64, color=BLUE, bold=True)
text(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.7),
     "1 aplikasi untuk semua peran · 0 kertas arsip manual · 100% tercatat & bisa diaudit — setiap rupiah, setiap galon, setiap perjalanan.",
     size=18, color=GRAY)
bullets(s, [
    "Langkah berikutnya: lengkapi data master (armada, driver, merk), pelatihan singkat per peran, dan prioritas pengembangan berikutnya.",
], y=Inches(4.4), size=17)
keybox(s, "🙏  Terima kasih. Siap menjawab pertanyaan dan mendampingi penerapan.")
footer(s, 12, 12)

OUT = "presentasi/BPF_WorkHub_Presentasi.pptx"
prs.save(OUT)
print(f"OK — {len(prs.slides.__iter__.__self__._sldIdLst)} slide → {OUT}")
